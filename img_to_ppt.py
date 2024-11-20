import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration

def ensure_output_folder(output_folder):
    """Ensure that the output folder exists and is empty."""
    if os.path.exists(output_folder):
        for filename in os.listdir(output_folder):
            file_path = os.path.join(output_folder, filename)
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)
    else:
        os.makedirs(output_folder)

def generate_text_for_image(image_path, processor, model):
    """Generate descriptive text for an image using BLIP."""
    with Image.open(image_path) as img:
        inputs = processor(img, return_tensors="pt")
        output = model.generate(**inputs)
        description = processor.decode(output[0], skip_special_tokens=True)
    return description

def add_image_and_text(slide, image_path, description, slide_width=10, slide_height=7.5, row_index=0, max_rows=3):
    """
    Adds a single image on the left and its corresponding description on the right.
    Positions each pair in a new row, with a maximum number of rows per slide.
    """
    left_column_width = Inches(slide_width / 2 - 0.5)
    right_column_x = Inches(slide_width / 2 + 0.5)
    row_height = (Inches(slide_height) - Inches(1)) / max_rows  # Height allocated per row

    # Position for image (left column)
    top = Inches(0.5) + row_index * row_height  # Adjust top position based on row index
    with Image.open(image_path) as img:
        aspect_ratio = img.width / img.height
        width = min(left_column_width, row_height * aspect_ratio)
        height = width / aspect_ratio
    slide.shapes.add_picture(image_path, Inches(0.5), top, width=width, height=height)

    # Position for text (right column) with font size adjustment and wrapping
    text_box = slide.shapes.add_textbox(right_column_x, top, left_column_width, row_height)
    text_frame = text_box.text_frame

    # Set initial font size and define maximum acceptable length
    max_text_length = 150
    initial_font_size = 14 if len(description) <= max_text_length else 12

    # Truncate text if it exceeds max length and adjust for ellipsis
    if len(description) > max_text_length:
        description = description[:max_text_length - 3] + "..."

    # Adjust font size to fit within bounds and add wrapped text
    text_frame.text = description
    text_frame.word_wrap = True  # Enable word wrapping
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(initial_font_size)

def create_ppt_with_images_and_text(base_folder_path, output_folder):
    ensure_output_folder(output_folder)
    output_ppt_path = os.path.join(output_folder, 'presentation_with_images_and_text.pptx')
    prs = Presentation()

    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    # Load BLIP processor and model for image captioning
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

    for folder_name in sorted(os.listdir(base_folder_path), key=lambda x: int(x.split('_')[1])):
        folder_path = os.path.join(base_folder_path, folder_name)
        if os.path.isdir(folder_path):
            images = [os.path.join(folder_path, f) for f in sorted(os.listdir(folder_path)) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
            if not images:
                print(f"Skipping empty folder: {folder_name}")
                continue

            for i in range(0, len(images), 3):  # Maximum 3 image-text pairs per slide
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                for row_index, image_path in enumerate(images[i:i+3]):
                    description = generate_text_for_image(image_path, processor, model)
                    add_image_and_text(slide, image_path, description, slide_width, slide_height, row_index)

    prs.save(output_ppt_path)
    print(f"PowerPoint presentation with images and descriptions saved at: {output_ppt_path}")

# Entry point of the script
if __name__ == "__main__":
    # Update these paths with your actual paths
    base_folder_path = 'output_visualizations'  # Path to the folder containing image folders (e.g., 'page_1', 'page_2', ...)
    output_folder = 'ppt'  # Path where the PowerPoint will be saved

    create_ppt_with_images_and_text(base_folder_path, output_folder)
